# -*- coding: utf-8 -*-
"""Build the AI Data Center Cooling sector-overview deck (Traditional Chinese)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ---- palette ----
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
BLUE   = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x2E, 0x9B, 0xD6)
TEAL   = RGBColor(0x14, 0xA3, 0x8C)
AMBER  = RGBColor(0xE8, 0x9A, 0x1C)
RED    = RGBColor(0xC0, 0x3A, 0x2B)
GREY   = RGBColor(0x5A, 0x5A, 0x5A)
LGREY  = RGBColor(0xEC, 0xEF, 0xF3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x22, 0x2A, 0x35)
CJK    = "Microsoft JhengHei"

EMUW, EMUH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width  = EMUW
prs.slide_height = EMUH
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line_sp=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.line_spacing = line_sp
        for (t, sz, b, c) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = CJK
    return tb

def header(s, kicker, title):
    rect(s, 0, 0, EMUW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), EMUW, Inches(0.06), ACCENT)
    txt(s, Inches(0.55), Inches(0.16), Inches(12), Inches(0.32),
        [[(kicker, 12, True, ACCENT)]])
    txt(s, Inches(0.55), Inches(0.46), Inches(12.3), Inches(0.62),
        [[(title, 25, True, WHITE)]])

def footer(s, n):
    txt(s, Inches(0.55), Inches(7.08), Inches(9), Inches(0.3),
        [[("AI 資料中心散熱 — 產業研究  |  資料時點 2026 年 6 月  |  僅供內部研究參考", 8, False, GREY)]])
    txt(s, Inches(12.4), Inches(7.08), Inches(0.7), Inches(0.3),
        [[(str(n), 9, False, GREY)]], align=PP_ALIGN.RIGHT)

def cell(c, text, size=10, bold=False, color=DARK, fill=None, align=PP_ALIGN.LEFT):
    c.text = ""
    tf = c.text_frame; tf.word_wrap = True
    c.margin_left = Pt(5); c.margin_right = Pt(4); c.margin_top = Pt(2); c.margin_bottom = Pt(2)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    for j, line in enumerate(text.split("\n")):
        if j > 0:
            p = tf.add_paragraph(); p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = CJK
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    else:
        c.fill.solid(); c.fill.fore_color.rgb = WHITE

def style_table(tbl):
    # diste default banding off
    try:
        tbl.first_row = True; tbl.horz_banding = False
    except Exception:
        pass

# =====================================================================
# 1. TITLE
# =====================================================================
s = slide()
rect(s, 0, 0, EMUW, EMUH, NAVY)
rect(s, 0, Inches(4.55), EMUW, Inches(0.06), ACCENT)
rect(s, Inches(0.55), Inches(1.5), Inches(0.16), Inches(2.6), ACCENT)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.5),
    [[("產業研究 SECTOR OVERVIEW", 15, True, ACCENT)]])
txt(s, Inches(0.9), Inches(1.95), Inches(11.8), Inches(1.9),
    [[("AI 資料中心散熱", 50, True, WHITE)],
     [("Data Center Cooling / Thermal Management", 22, False, RGBColor(0xBB,0xCD,0xE0))]], line_sp=1.05)
txt(s, Inches(0.9), Inches(4.85), Inches(11.6), Inches(1.2),
    [[("由「物理極限」強制驅動的結構性轉變 — 風冷在 ~20kW/機櫃觸頂，", 15, False, RGBColor(0xD7,0xE2,0xEE))],
     [("GB200 世代起所有機櫃級系統「非液冷不可」，由 2026 年約 $6,600–7,250 億雲端資本支出推動。", 15, False, RGBColor(0xD7,0xE2,0xEE))]], line_sp=1.15)
txt(s, Inches(0.9), Inches(6.55), Inches(11.6), Inches(0.5),
    [[("產業概覽 · 競爭格局 · 同業比較 · 標的清單      |      資料時點：2026 年 6 月", 12, True, RGBColor(0x9F,0xB6,0xCE))]])

# =====================================================================
# 2. EXECUTIVE SUMMARY
# =====================================================================
s = slide(); header(s, "EXECUTIVE SUMMARY", "投資重點：一個確定性極高的結構性轉變")
cards = [
    ("拐點已至", "2025 = 拐點年", "液冷從「尖端」變「基本盤」；NVIDIA GB200/GB300 參考架構直接把 D2C 液冷列為規格要求。", ACCENT),
    ("成長動能", "液冷 ~20% CAGR", "Dell'Oro：液冷 2025 ~$3B（年增約一倍）→ 2029 ~$7B；熱管理本世紀末規模追平 UPS。", TEAL),
    ("資金引擎", "2026 資本支出 ~$700B", "四大 + Oracle 合計年增 ~77%，約 75% 與 AI 相關；Goldman 估 FY25–30 累計 5.3 兆。", AMBER),
    ("選股主軸", "純度決定估值", "真正純散熱上市標的稀少：Vertiv、Auras、AVC、Envicool；倍數隨「散熱純度」走高。", BLUE),
]
x = Inches(0.55); cw = Inches(2.95); gap = Inches(0.18)
for (tag, big, body, col) in cards:
    rect(s, x, Inches(1.5), cw, Inches(3.0), WHITE, line=LGREY)
    rect(s, x, Inches(1.5), cw, Inches(0.12), col)
    txt(s, x+Inches(0.18), Inches(1.75), cw-Inches(0.36), Inches(0.4), [[(tag, 12, True, col)]])
    txt(s, x+Inches(0.18), Inches(2.18), cw-Inches(0.36), Inches(0.7), [[(big, 19, True, DARK)]], line_sp=1.0)
    txt(s, x+Inches(0.18), Inches(3.0), cw-Inches(0.36), Inches(1.4), [[(body, 11.5, False, GREY)]], line_sp=1.12)
    x = x + cw + gap

txt(s, Inches(0.55), Inches(4.8), Inches(12.3), Inches(1.9),
    [[("為何「非液冷不可」？", 14, True, NAVY)],
     [("• GPU 功耗階梯：H100 700W → B200 1,000W → GB300 1,400W → Vera Rubin ~1,800–2,300W → 2029 年預估 >4,000W（Dell'Oro）。", 12, False, DARK)],
     [("• 機櫃密度階梯：企業 ~15kW → H100 風冷 ~40kW → GB200 NVL72 ~120kW → Rubin Ultra 'Kyber' ~600kW（2027）→ Feynman ~1MW。", 12, False, DARK)],
     [("• 風冷 ~20kW/機櫃即觸頂；GB200 起每一代都超出風冷上限 2–15 倍 → 液冷成為強制規格，而非選配。", 12, True, BLUE)]], line_sp=1.2)
footer(s, 2)

# =====================================================================
# 3. MARKET SIZE (chart)
# =====================================================================
s = slide(); header(s, "MARKET SIZING", "市場規模：液冷是核心成長向量")
# left text
txt(s, Inches(0.55), Inches(1.45), Inches(4.7), Inches(4.9),
    [[("錨定數字（Dell'Oro 窄口徑）", 13, True, NAVY)],
     [("液冷 2025 約 $3B（年增約一倍），2029 達約 $7B。", 11.5, False, DARK)],
     [("資料中心實體基礎設施 (DCPI) 2030 突破 $80B，中雙位數 CAGR。", 11.5, False, DARK)],
     [("", 6, False, DARK)],
     [("廣口徑長天期（不可與上方混用）", 13, True, NAVY)],
     [("液冷 2031–2035 約 $18–29B，CAGR ~20%（Grand View / Mordor / Precedence）。", 11.5, False, DARK)],
     [("AI 專用液冷 2036 約 $17.8B，CAGR ~17%（IDTechEx / FMI）。", 11.5, False, DARK)],
     [("", 6, False, DARK)],
     [("⚠ 口徑差異約 2 倍 — 頭條用 Dell'Oro／Omdia，液冷深入用 IDTechEx／Mordor。", 11, True, RED)]], line_sp=1.18)
# chart: liquid cooling revenue
cd = CategoryChartData()
cd.categories = ["2023", "2024", "2025", "2027E", "2029E"]
cd.add_series("液冷市場 (US$ B, Dell'Oro 口徑)", (0.8, 1.5, 3.0, 5.0, 7.0))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(5.5), Inches(1.55), Inches(7.3), Inches(4.7), cd)
ch = gf.chart; ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False; ch.legend.font.size = Pt(10); ch.legend.font.name = CJK
plot = ch.plots[0]; plot.has_data_labels = True
dl = plot.data_labels; dl.number_format = '"$"0.0"B"'; dl.number_format_is_linked = False
dl.font.size = Pt(11); dl.font.bold = True; dl.font.name = CJK; dl.font.color.rgb = NAVY
ser = plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb = ACCENT
cat = ch.category_axis; cat.tick_labels.font.size = Pt(11); cat.tick_labels.font.name = CJK
ch.value_axis.has_major_gridlines = True
ch.value_axis.tick_labels.font.size = Pt(9); ch.value_axis.tick_labels.font.name = CJK
ch.has_title = True; ch.chart_title.text_frame.text = "液冷市場規模軌跡（Dell'Oro 製造商收入口徑）"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = CJK
footer(s, 3)

# =====================================================================
# 4. TECH SEGMENTATION
# =====================================================================
s = slide(); header(s, "TECHNOLOGY SEGMENTATION", "技術分段：Direct-to-Chip 冷板式是贏家")
rows = [
    ["技術", "份額 / 現況", "關鍵點"],
    ["風冷 (CRAC/CRAH)", "2024 約 54% 市場", "裝機主力，但 ~20kW/機櫃即達實用上限"],
    ["D2C 冷板式 ✓", "佔液冷市場 73–75%", "雲端建置主流；單相為主、雙相隨 TDP 滲透"],
    ["浸沒式", "2024 約佔液冷投資 17%", "利基；需機櫃密度穩定 >150–200kW 才放量"],
    ["後門熱交換器 RDHx", "83% 新建案一年內導入", "過渡/改造方案，常作為導入液冷的橋接"],
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.55), Inches(1.55), Inches(12.25), Inches(3.0)).table
tbl.columns[0].width = Inches(3.1); tbl.columns[1].width = Inches(3.2); tbl.columns[2].width = Inches(5.95)
for ri, row in enumerate(rows):
    tbl.rows[ri].height = Inches(0.62) if ri else Inches(0.5)
    for ci, val in enumerate(row):
        if ri == 0:
            cell(tbl.cell(ri, ci), val, 12, True, WHITE, NAVY)
        else:
            hi = rows[ri][0].endswith("✓")
            cell(tbl.cell(ri, ci), val, 11.5, ci == 0, DARK, RGBColor(0xE3,0xF1,0xF8) if hi else (LGREY if ri % 2 else WHITE))
style_table(tbl)
txt(s, Inches(0.55), Inches(4.85), Inches(12.3), Inches(1.9),
    [[("採用率與口徑提醒", 14, True, NAVY)],
     [("• DLC 企業採用率 2025 約 22%（與 2024 持平）；58% colo 業者回報客戶有液冷需求（Uptime Institute 2025）。", 12, False, DARK)],
     [("• 「70% 液冷 / 30% 風冷」是描述單一高密度 AI 新建案內部組合；「30% 液冷」才是整體市場收入占比（2027–28）。兩者勿混淆。", 12, False, DARK)],
     [("• 結論：D2C 冷板式吃下液冷主流；浸沒式短期維持利基，待密度再上一階。", 12, True, BLUE)]], line_sp=1.2)
footer(s, 4)

# =====================================================================
# 5. DEMAND DRIVERS (ladders)
# =====================================================================
s = slide(); header(s, "DEMAND DRIVERS", "需求驅動：功耗與密度的雙重階梯")
# GPU TDP ladder
txt(s, Inches(0.55), Inches(1.4), Inches(6), Inches(0.4), [[("GPU 功耗 (TDP) 階梯", 14, True, NAVY)]])
gpu = [("H100/H200", "700W", 0.32), ("B200", "1,000W", 0.45),
       ("GB300", "1,400W", 0.6), ("Vera Rubin", "1,800–2,300W", 0.82), ("2029E", ">4,000W", 1.0)]
y = Inches(1.9)
for (name, val, frac) in gpu:
    txt(s, Inches(0.55), y, Inches(1.9), Inches(0.35), [[(name, 11, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(2.5), y+Inches(0.04), Inches(int(4.2*frac*914400)), Inches(0.3), ACCENT)
    txt(s, Inches(2.5)+Inches(int(4.2*frac*914400))+Inches(0.08), y, Inches(1.3), Inches(0.35),
        [[(val, 11, True, BLUE)]], anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.52)
# Rack density ladder
txt(s, Inches(7.1), Inches(1.4), Inches(6), Inches(0.4), [[("機櫃功率密度階梯", 14, True, NAVY)]])
rack = [("企業傳統", "~15kW", 0.025), ("H100 風冷", "~40kW", 0.07),
        ("GB200 NVL72", "~120kW", 0.2), ("Rubin 'Kyber'", "~600kW", 0.6), ("Feynman", "~1MW", 1.0)]
y = Inches(1.9)
for (name, val, frac) in rack:
    txt(s, Inches(7.1), y, Inches(1.9), Inches(0.35), [[(name, 11, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(9.0), y+Inches(0.04), Inches(int(2.5*frac*914400))+Inches(0.05), Inches(0.3), TEAL)
    txt(s, Inches(9.0)+Inches(int(2.5*frac*914400))+Inches(0.1), y, Inches(1.4), Inches(0.35),
        [[(val, 11, True, BLUE)]], anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.52)
rect(s, Inches(0.55), Inches(4.95), Inches(12.25), Inches(1.55), LGREY)
txt(s, Inches(0.8), Inches(5.1), Inches(11.8), Inches(1.3),
    [[("四大驅動力", 13, True, NAVY)],
     [("① 資本支出引擎：2026 四大 + Oracle ~$660–725B（年增 ~77%），75% 與 AI 相關   ② PUE 連 6 年停滯（~1.54），液冷省 10–50% 能耗", 11.5, False, DARK)],
     [("③ 用水壓力：AI 資料中心 2028 用水估達 1.07 兆公升（約 11 倍）   ④ NVIDIA 參考架構直接把 D2C 液冷列為規格要求", 11.5, False, DARK)]], line_sp=1.18)
footer(s, 5)

# =====================================================================
# 6. VALUE CHAIN / LANDSCAPE
# =====================================================================
s = slide(); header(s, "COMPETITIVE LANDSCAPE", "競爭格局：價值鏈與領導者")
chain = [
    ("元件", "冷板 / 快接頭 / 歧管 / 泵", ["Boyd", "CoolIT→Ecolab", "AVC 3017", "Auras 3324", "Parker / Danfoss", "Lotes 3533 · Fositek 6805"], BLUE),
    ("CDU / 熱系統", "冷卻液分配與全套液冷", ["Vertiv VRT ★", "Schneider (Motivair)", "nVent NVT", "Boyd", "Envicool 002837"], ACCENT),
    ("伺服器 / 機櫃 ODM", "機櫃級整合與組裝", ["Foxconn 2317 ~40%", "Quanta 2382", "Wiwynn 6669", "Supermicro SMCI"], TEAL),
    ("使能者 / 服務", "熱交換器 · 快接 · 部署服務", ["Kaori 8210/8996", "Fositek 6805", "Vertiv 服務", "浸沒：GRC·Submer·LiquidStack"], AMBER),
]
x = Inches(0.55); cw = Inches(3.0); gap = Inches(0.12)
for (title, sub, names, col) in chain:
    rect(s, x, Inches(1.5), cw, Inches(0.95), col)
    txt(s, x+Inches(0.15), Inches(1.6), cw-Inches(0.3), Inches(0.45), [[(title, 15, True, WHITE)]])
    txt(s, x+Inches(0.15), Inches(2.02), cw-Inches(0.3), Inches(0.4), [[(sub, 10, False, RGBColor(0xEA,0xF2,0xF8))]])
    rect(s, x, Inches(2.5), cw, Inches(3.6), WHITE, line=LGREY)
    yy = Inches(2.65)
    for nm in names:
        rect(s, x+Inches(0.12), yy, cw-Inches(0.24), Inches(0.52), LGREY)
        txt(s, x+Inches(0.25), yy, cw-Inches(0.4), Inches(0.52), [[(nm, 11, "★" in nm, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + Inches(0.62)
    if title != "使能者 / 服務":
        txt(s, x+cw-Inches(0.02), Inches(3.6), Inches(0.18), Inches(0.5), [[("›", 22, True, GREY)]])
    x = x + cw + gap
txt(s, Inches(0.55), Inches(6.25), Inches(12.3), Inches(0.7),
    [[("誰在搶份額：", 12, True, NAVY), ("Vertiv（份額第一）· nVent（backlog 約 3× 最突出新晉）· Schneider（整併中的老二）· 台灣 AVC/Auras（純散熱直接受惠）。  ", 11.5, False, DARK),
      ("警示：", 12, True, RED), ("Asetek（已淡出 DC）· Parker（散熱僅 ~1%）。", 11.5, False, DARK)]], line_sp=1.15)
footer(s, 6)

# =====================================================================
# 7. COMPS TABLE
# =====================================================================
s = slide(); header(s, "PEER COMPARISON", "同業比較 (Comps)  —  資料約 2026/6/17–18")
comps = [
    ["公司", "代碼", "市值(USD)", "營收YoY", "營益率", "Fwd P/E", "EV/EBITDA", "1年報酬", "純度"],
    ["Vertiv", "VRT", "~$122B", "+28%", "18.8%", "~46x", "~51x", "+173%", "純(高)"],
    ["nVent", "NVT", "~$28B", "+11%", "16.2%", "~35x", "~32x", "+143%", "多元"],
    ["Schneider", "SU.PA", "~$160B", "+5%", "17.4%", "~28x", "~21x", "+27%", "多元"],
    ["Supermicro", "SMCI", "~$18B", "+123%", "4.5%", "~9x", "~16x", "−36%", "伺服器OEM"],
    ["Delta 台達電", "2308", "~$187B", "+33%", "17.8%", "~46x", "n/a", "+454%", "多元"],
    ["Auras 奇鋐", "3324", "~$3.2B", "+48%", "14.1%", "~21x", "n/a", "+132%", "純散熱"],
    ["AVC", "3017", "~$31B", "+107%", "n/a", "~23x", "n/a", "+217%", "DC散熱"],
    ["Wiwynn 緯穎", "6669", "~$30B", "+129%", "6.6%", "~13.5x", "~13.6x", "+103%", "ODM"],
    ["Lotes 嘉澤", "3533", "~$8.2B", "+12%", "30.0%", "~22.8x", "~18.5x", "+67%", "連接器"],
    ["Envicool 英維克", "002837", "~$15B", "+32%", "8.6%", "~77x", "n/a", "+237%", "純(最近)"],
    ["Parker ⚠", "PH", "~$119B", "+6%", "21.7%", "~28x", "~23x", "+43%", "~1%散熱"],
]
tbl = s.shapes.add_table(len(comps), 9, Inches(0.4), Inches(1.45), Inches(12.55), Inches(4.7)).table
widths = [1.75, 0.95, 1.35, 1.05, 1.0, 1.05, 1.3, 1.1, 2.0]
for i, w in enumerate(widths):
    tbl.columns[i].width = Inches(w)
for ri, row in enumerate(comps):
    tbl.rows[ri].height = Inches(0.34)
    for ci, val in enumerate(row):
        if ri == 0:
            cell(tbl.cell(ri, ci), val, 10, True, WHITE, NAVY, PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT)
        else:
            col = DARK
            if ci == 7 and val.startswith("−"): col = RED
            elif ci == 7: col = TEAL
            fill = LGREY if ri % 2 else WHITE
            cell(tbl.cell(ri, ci), val, 9.5, ci == 0, col, fill, PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT)
style_table(tbl)
txt(s, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.7),
    [[("讀法：", 11, True, NAVY), ("① 估值隨「散熱純度」走高（Vertiv ~51x vs. Schneider ~21x）；② 元件 vs. 系統的 EV/Rev 不可直接比（ODM 為低毛利轉手）；③ 台股中小型 EV 倍數免費源缺漏，標 n/a。代碼校正：Auras=3324、AVC=3017。", 10, False, GREY)]], line_sp=1.12)
footer(s, 7)

# =====================================================================
# 8. WATCHLIST
# =====================================================================
s = slide(); header(s, "WATCHLIST", "值得關注標的清單")
cols = [
    ("核心持有", TEAL, "最乾淨的曝險", [
        ("Vertiv (VRT)", "AI 散熱領頭羊；端到端 + NVIDIA 綁定；backlog >$15B"),
        ("Auras 奇鋐 (3324)", "台灣純散熱高 beta；2026 指引 +60–70%；CDU 放量"),
        ("AVC (3017)", "台灣熱模組龍頭；BOM 最廣；能見度看到 2029"),
    ]),
    ("強力觀察", ACCENT, "搶份額 / 轉機", [
        ("nVent (NVT)", "最突出新晉搶份額；backlog 約 3×；估值較合理"),
        ("Schneider (SU.PA)", "整併中的老二(Motivair)；grid-to-chip；大型股最便宜"),
        ("Delta 台達電 (2308)", "電源+散熱一體；切入 800VDC AI 工廠"),
        ("Kaori 高力 (8996)", "GB200 板式熱交換器隱形冠軍"),
        ("Fositek/Lotes (6805/3533)", "快接頭/連接器卡點；Rubin 世代受惠"),
    ]),
    ("注意風險", RED, "瑕疵 / 不純", [
        ("Supermicro (SMCI)", "DLC 領先但毛利薄、治理爭議；唯一負報酬"),
        ("Envicool (002837)", "最接近純散熱，但 Q1 淨利 −82%、估值高"),
        ("Parker / Asetek / INVT", "散熱曝險過小或業務休眠 — 非純標的"),
    ]),
]
x = Inches(0.55); cw = Inches(4.0); gap = Inches(0.1)
for (title, col, sub, items) in cols:
    rect(s, x, Inches(1.5), cw, Inches(0.85), col)
    txt(s, x+Inches(0.18), Inches(1.58), cw-Inches(0.3), Inches(0.45), [[(title, 16, True, WHITE)]])
    txt(s, x+Inches(0.18), Inches(2.0), cw-Inches(0.3), Inches(0.35), [[(sub, 11, False, WHITE)]])
    rect(s, x, Inches(2.4), cw, Inches(4.35), WHITE, line=LGREY)
    yy = Inches(2.55)
    for (nm, desc) in items:
        txt(s, x+Inches(0.15), yy, cw-Inches(0.3), Inches(0.32), [[(nm, 12, True, NAVY)]])
        txt(s, x+Inches(0.15), yy+Inches(0.30), cw-Inches(0.3), Inches(0.55), [[(desc, 9.8, False, GREY)]], line_sp=1.05)
        yy = yy + Inches(0.83)
    x = x + cw + gap
footer(s, 8)

# =====================================================================
# 9. RISKS + SOURCES
# =====================================================================
s = slide(); header(s, "RISKS & DISCLOSURES", "風險、方法論與來源")
txt(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(5.0),
    [[("關鍵風險", 14, True, NAVY)],
     [("• 改造經濟性：改造比新建貴 50–100%，每 MW $2–3M，週期 18–36 月", 11, False, DARK)],
     [("• 缺乏標準：~40 家做 CDU、品質不一；漏液佔事故 19% 且上升", 11, False, DARK)],
     [("• 整併：Dell'Oro 估本世紀末存活 <10 家拿走絕大多數份額", 11, False, DARK)],
     [("• 電網瓶頸：美國併網佇列卡關 ~2,300 GW；電力已取代硬體成瓶頸", 11, False, DARK)],
     [("• 消化風險：資本支出(~$700B) vs. 直接 AI 營收(~$51B)落差，2027–28 疑慮", 11, False, DARK)],
     [("", 6, False, DARK)],
     [("方法論注意", 14, True, NAVY)],
     [("• 市場規模勿跨機構混用（口徑差約 2 倍）；頭條用 Dell'Oro／Omdia", 11, False, DARK)],
     [("• 微流道、Rubin 600kW/1MW 機櫃屬路線圖/原型，非已出貨規格", 11, False, DARK)],
     [("• 台股/中國中小型 EV 倍數與散熱營收占比多為估計值", 11, False, DARK)]], line_sp=1.16)
txt(s, Inches(6.9), Inches(1.45), Inches(5.9), Inches(5.0),
    [[("主要來源", 14, True, NAVY)],
     [("• Dell'Oro Group — DCPI / 液冷市場 / GPU TDP", 11, False, DARK)],
     [("• Omdia — 市場規模與 CAGR", 11, False, DARK)],
     [("• IDTechEx — 技術分段份額（D2C ~75%、浸沒 ~17%）", 11, False, DARK)],
     [("• Uptime Institute 2025 — DLC 採用率、PUE", 11, False, DARK)],
     [("• JLL / McKinsey / SemiAnalysis — 需求與機櫃密度", 11, False, DARK)],
     [("• NVIDIA GTC/CES — 路線圖（GB300 / Rubin / Feynman）", 11, False, DARK)],
     [("• 公司 IR / SEC 8-K / Digitimes / TrendForce — 公司財務", 11, False, DARK)],
     [("", 10, False, DARK)],
     [("免責聲明", 13, True, RED)],
     [("本資料彙整自第三方研究與公開資訊，僅供內部研究參考，不構成投資建議。估值與股價為即時變動數據，下單前請核對主要申報文件與中文公司名/代碼（3324 奇鋐 vs 3017 AVC）。", 10, False, GREY)]], line_sp=1.16)
footer(s, 9)

prs.save("./out/AI資料中心散熱_產業研究.pptx")
print("SAVED ./out/AI資料中心散熱_產業研究.pptx  slides:", len(prs.slides._sldIdLst))
