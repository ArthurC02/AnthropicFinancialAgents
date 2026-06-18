# -*- coding: utf-8 -*-
"""Vertiv (VRT) M&A proposal deck (first draft)."""
import os
os.makedirs("./out", exist_ok=True)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x0B,0x1F,0x3A); BLUE=RGBColor(0x1F,0x4E,0x79); ACCENT=RGBColor(0x2E,0x9B,0xD6)
TEAL=RGBColor(0x14,0xA3,0x8C); AMBER=RGBColor(0xE8,0x9A,0x1C); RED=RGBColor(0xC0,0x3A,0x2B)
WHITE=RGBColor(0xFF,0xFF,0xFF); GREY=RGBColor(0x5A,0x5A,0x5A); LGREY=RGBColor(0xEC,0xEF,0xF3); DARK=RGBColor(0x22,0x2A,0x35)
CJK="Microsoft JhengHei"
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; EMW=Inches(13.333)

def rect(s,x,y,w,h,fill,line=None):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(0.75)
    sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=4,line_sp=1.0):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.line_spacing=line_sp
        for (t,sz,b,c) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=CJK
    return tb
def header(s,kicker,title,n):
    rect(s,0,0,EMW,Inches(1.1),NAVY); rect(s,0,Inches(1.1),EMW,Inches(0.055),ACCENT)
    txt(s,Inches(0.55),Inches(0.16),Inches(12),Inches(0.3),[[(kicker,12,True,ACCENT)]])
    txt(s,Inches(0.55),Inches(0.45),Inches(12.3),Inches(0.6),[[(title,24,True,WHITE)]])
    txt(s,Inches(0.55),Inches(7.06),Inches(11),Inches(0.3),[[("Vertiv (VRT) 併購提案初稿 · 僅供內部討論 · 示意性,非投資建議",8,False,GREY)]])
    txt(s,Inches(12.5),Inches(7.06),Inches(0.6),Inches(0.3),[[(str(n),9,False,GREY)]],align=PP_ALIGN.RIGHT)
def table(s,x,y,rows,colw,rowh=Inches(0.36),fs=10,header_fill=NAVY):
    nr=len(rows); nc=len(rows[0])
    tb=s.shapes.add_table(nr,nc,x,y,sum(colw,Inches(0)),rowh*nr).table
    for ci,w in enumerate(colw): tb.columns[ci].width=w
    for ri,row in enumerate(rows):
        tb.rows[ri].height=rowh
        for ci,val in enumerate(row):
            c=tb.cell(ri,ci); c.text=""; c.margin_left=Pt(5); c.margin_right=Pt(4); c.margin_top=Pt(1); c.margin_bottom=Pt(1)
            c.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=str(val); r.font.name=CJK; r.font.size=Pt(fs if ri else fs)
            if ri==0: r.font.bold=True; r.font.color.rgb=WHITE; c.fill.solid(); c.fill.fore_color.rgb=header_fill
            else:
                r.font.color.rgb=DARK; r.font.bold=(ci==0)
                c.fill.solid(); c.fill.fore_color.rgb=(LGREY if ri%2 else WHITE)
    return tb

# 1 COVER
s=prs.slides.add_slide(BLANK); rect(s,0,0,EMW,Inches(7.5),NAVY); rect(s,0,Inches(4.5),EMW,Inches(0.06),ACCENT)
rect(s,Inches(0.55),Inches(1.4),Inches(0.16),Inches(2.5),ACCENT)
txt(s,Inches(0.9),Inches(1.3),Inches(11),Inches(0.5),[[("併購提案初稿 M&A PROPOSAL — DRAFT",15,True,ACCENT)]])
txt(s,Inches(0.9),Inches(1.95),Inches(11.6),Inches(1.6),
    [[("收購 Vertiv (VRT) 之評估",40,True,WHITE)],[("Project Coolwater — 策略與財務買方視角",20,False,RGBColor(0xBB,0xCD,0xE0))]],line_sp=1.05)
txt(s,Inches(0.9),Inches(4.75),Inches(11.6),Inches(1.6),
    [[("同業比較 · 先例交易 · DCF · 示意 LBO · 足球場估值區間",14,True,RGBColor(0xD7,0xE2,0xEE))],
     [("核心結論:VRT 為 AI 散熱龍頭,但現價已 price in 完美成長;內在/交易法公允值 $158–$245(DCF 低端 $100)vs 現價 $320,",12,False,RGBColor(0xC9,0xD8,0xE8))],
     [("控制溢價難以在基本面上 justify,且傳統 LBO 因規模/倍數不可行。",12,False,RGBColor(0xC9,0xD8,0xE8))]],line_sp=1.15)
txt(s,Inches(0.9),Inches(6.7),Inches(11.6),Inches(0.4),[[("資料時點 2026 年 6 月 · 僅供內部討論,非投資建議",11,True,RGBColor(0x9F,0xB6,0xCE))]])

# 2 FRAMING
s=prs.slides.add_slide(BLANK); header(s,"FRAMING & DISCLAIMER","框架假設與重要聲明",2)
txt(s,Inches(0.55),Inches(1.45),Inches(12.3),Inches(5.2),
    [[("本提案的框架假設",15,True,NAVY)],
     [("• 視角:以「潛在收購方評估收購 VRT」為框架。DCF/Comps/先例交易為策略買方視角;示意 LBO 為財務買方(PE)視角。",12,False,DARK)],
     [("• VRT 現況:市值 ~$122bn、企業價值 ~$124bn、FY26E EV/EBITDA ~40x(trailing ~51x)——估值已屬全球同業最高。",12,False,DARK)],
     [("• 示意控制報價設為 $400/股(對現價 $320 之 +25% 溢價),用於示範交易與報酬數學。",12,False,DARK)],
     [("",6,False,DARK)],
     [("重要聲明",15,True,RED)],
     [("• 本檔為初稿,所有財務數字為公開資訊與第三方彙整之估計值(現價、股數、淨負債、EBITDA、同業/交易倍數),正式提案前須以 FactSet/Bloomberg 與公司申報文件校正。",12,False,DARK)],
     [("• 足球場圖中『52週區間』與『分析師目標價』為示意值;LBO 因 VRT 規模 ~$122bn 在實務上不可行,僅作報酬數學示範。",12,False,DARK)],
     [("• 本資料不構成投資建議或要約。",12,True,BLUE)]],line_sp=1.18)

# 3 EXEC SUMMARY
s=prs.slides.add_slide(BLANK); header(s,"EXECUTIVE SUMMARY","執行摘要:策略誘人,但估值是硬傷",3)
cards=[("策略理由 ✓","AI 散熱龍頭","端到端組合 + NVIDIA 參考設計綁定;backlog >$15bn;液冷 ~40% CAGR。稀缺的純標的。",TEAL),
 ("估值缺口 ✗","公允值 < 現價","DCF/先例/LBO 指向 $158–$245(DCF 低端 $100),低於現價 $320;再加控制溢價更難 justify。",RED),
 ("LBO 不可行 ✗","規模/倍數障礙","~$122bn 規模 + 40x EBITDA,槓桿空間極小;示意 $400 報價 IRR ~0–1%。",AMBER),
 ("建議方向","結構性替代","直接全現金收購報酬不佳;宜考慮少數股權/合資/可轉換結構,或等待更佳進場點。",BLUE)]
x=Inches(0.55); cw=Inches(2.95)
for tag,big,body,col in cards:
    rect(s,x,Inches(1.45),cw,Inches(3.1),WHITE,line=LGREY); rect(s,x,Inches(1.45),cw,Inches(0.12),col)
    txt(s,x+Inches(0.18),Inches(1.68),cw-Inches(0.36),Inches(0.4),[[(tag,12,True,col)]])
    txt(s,x+Inches(0.18),Inches(2.12),cw-Inches(0.36),Inches(0.7),[[(big,17,True,DARK)]],line_sp=1.0)
    txt(s,x+Inches(0.18),Inches(2.95),cw-Inches(0.36),Inches(1.5),[[(body,11,False,GREY)]],line_sp=1.12)
    x=x+cw+Inches(0.18)
txt(s,Inches(0.55),Inches(4.85),Inches(12.3),Inches(1.9),
    [[("一句話建議",14,True,NAVY)],
     [("VRT 是同題材中品質最高的資產,但市場定價已充分(甚至過度)反映其成長。以併購角度,在 $320 之上再付控制溢價,",12,False,DARK)],
     [("無法以 DCF、先例交易或 LBO 報酬支撐。建議:① 暫不推進全額收購;② 探索結構性/少數股權合作;③ 將 VRT 列入觀察,等待估值修正後的進場窗口。",12,True,BLUE)]],line_sp=1.2)

# 4 VRT SNAPSHOT
s=prs.slides.add_slide(BLANK); header(s,"TARGET SNAPSHOT","標的快照:Vertiv (VRT)",4)
table(s,Inches(0.55),Inches(1.5),
 [["指標","數值","說明"],
  ["FY25 營收","~$10.84bn (+28% YoY)","AI 資料中心基礎設施(電力+散熱+服務)"],
  ["FY26E 營收(指引中點)","~$13.75bn","年增約 27%"],
  ["EBITDA 利潤率","~22–23%","營運槓桿持續擴張"],
  ["Backlog","> $15bn","翻倍以上;訂單能見度高"],
  ["市值 / 企業價值","~$122bn / ~$124bn","股數 ~0.382bn,淨負債 ~$1.9bn"],
  ["EV/EBITDA","~40x fwd(~51x trailing)","全球同業最高,反映 AI 純度與龍頭地位"],
  ["護城河","端到端 + NVIDIA 綁定","CoolTera/PurgeRite 等併購擴張液冷"]],
 [Inches(3.4),Inches(3.6),Inches(5.25)],rowh=Inches(0.5),fs=11)
txt(s,Inches(0.55),Inches(6.35),Inches(12.3),Inches(0.7),
    [[("為何是標的:",12,True,NAVY),("AI 散熱最乾淨的純標的、龍頭份額、與 NVIDIA 路線圖綁定——對欲切入 AI 基礎設施的策略買方具高度吸引力。",11,False,DARK)]],line_sp=1.1)

# 5 FOOTBALL FIELD (money slide)
s=prs.slides.add_slide(BLANK); header(s,"VALUATION SUMMARY","估值區間總覽:足球場圖",5)
s.shapes.add_picture("./out/charts/vrt_football.png",Inches(0.5),Inches(1.35),width=Inches(8.6))
txt(s,Inches(9.3),Inches(1.5),Inches(3.7),Inches(5.2),
    [[("讀法",14,True,NAVY)],
     [("• 內在(DCF)、先例交易、LBO ability-to-pay 三法群聚 $158–$245(DCF 低端達 $100)。",11,False,DARK)],
     [("• 現價 $320 已在多數方法之上。",11,True,RED)],
     [("• 控制溢價(→$400)僅情緒類參考(52週高、樂觀PT)可及。",11,False,DARK)],
     [("",6,False,DARK)],
     [("意涵",13,True,NAVY)],
     [("收購 VRT = 對『已偏貴的價格』再付溢價;除非有龐大綜效或極長期 AI 跑道,否則難以創造收購方價值。",11,True,BLUE)]],line_sp=1.16)

# 6 COMPS
s=prs.slides.add_slide(BLANK); header(s,"TRADING COMPS","同業比較",6)
table(s,Inches(0.55),Inches(1.5),
 [["公司","代碼","市值$bn","營收YoY","EV/EBITDA","Fwd P/E"],
  ["Vertiv","VRT","~122","+28%","~40–51x","~46x"],
  ["nVent","NVT","~28","+11%","~32x","~35x"],
  ["Schneider","SU.PA","~160","+5%","~21x","~28x"],
  ["Delta 台達電","2308","~187","+33%","n/a","~46x"],
  ["Supermicro","SMCI","~18","+123%","~16x","~9x"],
  ["同業中位(參考)","—","—","—","~26x","~30x"]],
 [Inches(2.6),Inches(1.5),Inches(1.9),Inches(1.9),Inches(2.3),Inches(2.05)],rowh=Inches(0.42),fs=11)
txt(s,Inches(0.55),Inches(5.2),Inches(12.3),Inches(1.4),
    [[("Comps 隱含區間",13,True,NAVY)],
     [("套用同業 25–40x 於 VRT FY26E EBITDA ~$3.1bn → 每股約 $198–$320。",12,False,DARK)],
     [("注意:VRT 本身即同業最高倍數;以 comps 反推時其『區間頂端』本質上就是現價,故 comps 對收購溢價的支撐有限。",12,True,BLUE)]],line_sp=1.18)

# 7 PRECEDENTS
s=prs.slides.add_slide(BLANK); header(s,"PRECEDENT TRANSACTIONS","先例交易",7)
table(s,Inches(0.55),Inches(1.45),
 [["日期","買方","標的","EV","EV/EBITDA"],
  ["2026-03","Ecolab","CoolIT","$4.75bn","29x NTM / 24x'27E"],
  ["2025-12","Eaton","Boyd Thermal","$9.5bn","22.5x '26E"],
  ["2024-10","Schneider","Motivair (75%)","$0.85bn","中個位數(rev)"],
  ["2025-11","Vertiv","PurgeRite","~$1.0bn","~10x(含綜效)"],
  ["2023-04","Carrier","Viessmann(參考)","~$13bn","~13x"],
  ["2025-03","nVent","Avail EPG(電力)","$0.975bn","~2.6x(rev)"]],
 [Inches(1.5),Inches(1.9),Inches(3.2),Inches(2.0),Inches(3.65)],rowh=Inches(0.42),fs=10.5)
txt(s,Inches(0.55),Inches(5.25),Inches(12.3),Inches(1.4),
    [[("先例隱含區間",13,True,NAVY)],
     [("純散熱控制交易 EV/EBITDA ~10–29x(前瞻);套 20–28x 於 VRT FY26E EBITDA → 每股約 $158–$223。",12,False,DARK)],
     [("注意:CoolIT/Boyd 為前瞻倍數且標的更小、成長更快;VRT 規模大、成長較緩,直接類比會高估。資料為第三方彙整,須查核。",12,True,BLUE)]],line_sp=1.18)

# 8 DCF
s=prs.slides.add_slide(BLANK); header(s,"DCF","DCF 估值摘要",8)
table(s,Inches(0.55),Inches(1.5),
 [["方法","WACC","終值","隱含每股","備註"],
  ["永續成長法","11.4%","g=3.5%","~$102","紀律性;終值占EV偏高"],
  ["出場倍數法","11.4%","22x EBITDA","~$245","較寬鬆"],
  ["DCF 區間","—","—","$100–$245","中樞遠低於現價 $320"]],
 [Inches(2.6),Inches(1.7),Inches(2.4),Inches(2.4),Inches(3.55)],rowh=Inches(0.5),fs=11)
txt(s,Inches(0.55),Inches(4.2),Inches(12.3),Inches(2.4),
    [[("關鍵",14,True,NAVY)],
     [("• WACC 11.4% 反映 VRT 高 beta(~1.45);5 年顯性期 + 永續法對高成長股偏保守。",12,False,DARK)],
     [("• 即便用寬鬆的 22x 出場倍數,DCF 上限 ~$245 仍低於現價 $320。",12,False,DARK)],
     [("• 結論:以紀律性 DCF,VRT 內在價值中樞 $100–$245,顯著低於現價——市場已 price in 高度完美成長。",12,True,BLUE)],
     [("• 完整模型(含三情境/敏感度)見隨附 Vertiv_VRT_DCF.xlsx。",11,False,GREY)]],line_sp=1.2)

# 9 LBO
s=prs.slides.add_slide(BLANK); header(s,"ILLUSTRATIVE LBO","示意 LBO 與 Ability-to-Pay",9)
table(s,Inches(0.55),Inches(1.5),
 [["示意 LBO(報價 $400 / +25%)","值"],
  ["進場 EV / EV-EBITDA","~$155bn / ~50x"],
  ["新負債(6x EBITDA)","~$18.6bn"],
  ["發起人股權(plug)","~$140bn"],
  ["出場(FY31,22x)股權","~$145bn"],
  ["MOIC / IRR","1.04x / ~0.7%"]],
 [Inches(4.4),Inches(2.2)],rowh=Inches(0.44),fs=11)
table(s,Inches(7.4),Inches(1.5),
 [["Ability-to-Pay","每股 $"],
  ["15% IRR","~$233"],
  ["20% IRR","~$196"],
  ["25% IRR","~$168"],
  ["現價(對照)","$320"]],
 [Inches(3.3),Inches(2.0)],rowh=Inches(0.44),fs=11)
txt(s,Inches(0.55),Inches(5.2),Inches(12.3),Inches(1.5),
    [[("解讀(含綜效情境)",13,True,NAVY)],
     [("• 示意 $400 報價 IRR ~0–1%(進場 ~50x、槓桿空間極小、發起人需寫 ~$140bn 股權)。要達 20–25% IRR 進場價需 ~$168–$196。",12,False,DARK)],
     [("• 綜效情境(@20% IRR):綜效 $0 → $196;$1.0bn → $220;$2.0bn → $243——即便 $2bn 綜效(≈進場 EBITDA 64%)仍 < 現價 $320。",12,False,DARK)],
     [("• 反推:要 justify $320 需 run-rate 綜效 ~$5.3bn(不切實際)。完整綜效表與 5×5 IRR 敏感度見工作簿 Synergies 分頁。",12,True,RED)]],line_sp=1.16)

# 10 IRR HEATMAP
s=prs.slides.add_slide(BLANK); header(s,"LBO SENSITIVITY","LBO IRR 敏感度:報價 × 出場倍數",10)
s.shapes.add_picture("./out/charts/vrt_irr_heatmap.png",Inches(0.6),Inches(1.4),width=Inches(8.4))
txt(s,Inches(9.2),Inches(1.5),Inches(3.8),Inches(5.2),
    [[("讀法",14,True,NAVY)],
     [("• 縱軸=報價/股,橫軸=出場 EV/EBITDA;格內為 5 年 IRR。",11,False,DARK)],
     [("• 基準($400 / 22x,框線)= ~0.7% IRR。",11,True,RED)],
     [("• 整張表無任一組合達 PE 門檻(~20–25%)。",11,False,DARK)],
     [("• 即便報價降至 $320 + 出場 26x,IRR 也僅 ~10%。",11,False,DARK)],
     [("",6,False,DARK)],
     [("意涵",13,True,NAVY)],
     [("LBO 報酬對『進場價』高度敏感,但 VRT 在現價附近的進場根本無法產生 PE 等級報酬——財務買方路徑不通。",11,True,BLUE)]],line_sp=1.16)

# 11 CONCLUSION
s=prs.slides.add_slide(BLANK); header(s,"RECOMMENDATION & NEXT STEPS","結論、風險與後續",11)
txt(s,Inches(0.55),Inches(1.45),Inches(6.1),Inches(5.2),
    [[("建議",14,True,NAVY)],
     [("① 暫不推進對 VRT 的全額現金收購——估值缺口過大。",12,False,DARK)],
     [("② 若策略上必須布局,優先考慮:少數股權、合資、商業合作或可轉換結構,降低進場估值風險。",12,False,DARK)],
     [("③ 將 VRT 列入觀察名單,設定估值修正後的進場觸發(如 EV/EBITDA 回落至 ~25–30x)。",12,False,DARK)],
     [("④ 平行評估較小、估值合理的純標的(如台股 Auras 3324、AVC 3017,或私有液冷資產)作為替代切入。",12,False,DARK)]],line_sp=1.2)
txt(s,Inches(6.9),Inches(1.45),Inches(5.9),Inches(5.2),
    [[("風險與後續工作",14,True,NAVY)],
     [("• 數據校正:現價/股數/淨負債/EBITDA/同業與交易倍數以 FactSet/Bloomberg 與申報文件核實。",12,False,DARK)],
     [("• 綜效量化:策略買方若有明確成本/營收綜效,可上修 ability-to-pay——下一步應建綜效情境。",12,False,DARK)],
     [("• 反向風險:若 AI 需求與液冷滲透超預期,VRT 成長跑道可能比 5 年 DCF 更長,內在價值上修。",12,False,DARK)],
     [("• 後續:① 校正數據 ② 建綜效情境 ③ 結構性方案設計 ④ 與內部投委會討論進場觸發。",12,True,BLUE)]],line_sp=1.2)
rect(s,Inches(0.55),Inches(6.55),Inches(12.25),Inches(0.5),LGREY)
txt(s,Inches(0.7),Inches(6.6),Inches(12),Inches(0.4),
    [[("免責:初稿,示意性,非投資建議或要約;財務數字為估計值,須獨立查核。",10,True,GREY)]])

prs.save("./out/VRT_M&A_Proposal_Draft.pptx")
import os as _os
_os.remove("./out/deck_content.md") if _os.path.exists("./out/deck_content.md") else None
print("SAVED ./out/VRT_M&A_Proposal_Draft.pptx slides:",len(prs.slides._sldIdLst))
